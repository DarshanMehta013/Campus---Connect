import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions: 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    COLOR_PRIMARY = RGBColor(15, 23, 42)      # Deep Navy/Slate 900
    COLOR_SECONDARY = RGBColor(30, 41, 59)   # Slate 800
    COLOR_ACCENT = RGBColor(37, 99, 235)      # Blue 600
    COLOR_ACCENT_LIGHT = RGBColor(59, 130, 246) # Blue 500
    COLOR_TEAL = RGBColor(13, 148, 136)       # Teal 600
    COLOR_EMERALD = RGBColor(16, 185, 129)   # Emerald 500
    COLOR_AMBER = RGBColor(245, 158, 11)     # Amber 500
    COLOR_RED = RGBColor(239, 68, 68)        # Red 500
    COLOR_BG_LIGHT = RGBColor(248, 250, 252) # Slate 50
    COLOR_CARD_BG = RGBColor(255, 255, 255)  # Pure White
    COLOR_CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)   # Slate 900
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    COLOR_TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def add_header(slide, title_text, category_text="SEMINAR-II PRESENTATION | DIPLOMA IN COMPUTER ENGINEERING"):
        # Top banner category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = 'Calibri'

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.65))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = 'Segoe UI'

        # Underline accent bar
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.38), Inches(1.8), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.color.rgb = COLOR_ACCENT

    def set_slide_background(slide, color=COLOR_BG_LIGHT):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_PRIMARY)

    # Decorative accent card
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11.0), Inches(4.5))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "DIPLOMA IN COMPUTER ENGINEERING — SEMINAR-II"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_LIGHT
    p0.font.name = 'Calibri'
    p0.space_after = Pt(12)

    p1 = tf1.add_paragraph()
    p1.text = "CAMPUS CONNECT"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = 'Segoe UI'

    p2 = tf1.add_paragraph()
    p2.text = "Automated Multi-Role Resolution & Campus Governance Portal"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_EMERALD
    p2.font.name = 'Segoe UI'
    p2.space_after = Pt(24)

    p3 = tf1.add_paragraph()
    p3.text = "A Closed-Loop Digital Governance System: Student → Admin → Technician → Faculty → Resolution"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_LIGHT
    p3.font.name = 'Calibri'

    # Project Metadata Cards
    meta_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.0), Inches(10.8), Inches(1.8))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = COLOR_SECONDARY
    meta_box.line.color.rgb = RGBColor(51, 65, 85)

    mtf = meta_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.4)
    mtf.margin_top = Inches(0.25)

    mp1 = mtf.paragraphs[0]
    mp1.text = "PROJECT DETAILS & CANDIDATE INFORMATION"
    mp1.font.size = Pt(11)
    mp1.font.bold = True
    mp1.font.color.rgb = COLOR_ACCENT_LIGHT

    mp2 = mtf.add_paragraph()
    mp2.text = "• Domain: Web Engineering & Smart Governance     • Technology Stack: HTML, CSS, JavaScript, Bootstrap, PHP, MySQL"
    mp2.font.size = Pt(12)
    mp2.font.color.rgb = COLOR_TEXT_LIGHT

    mp3 = mtf.add_paragraph()
    mp3.text = "• Target Roles: Students | Admin Office | Specialized Technicians | Department Faculty Advisors"
    mp3.font.size = Pt(12)
    mp3.font.color.rgb = COLOR_TEXT_LIGHT

    add_speaker_notes(s1, 
        "Good morning respected evaluators, project guide, and dear fellow students.\n"
        "I am presenting my Seminar-II project titled 'Campus Connect: Automated Multi-Role Resolution Portal'.\n"
        "This project is developed for the Diploma in Computer Engineering curriculum under the Smart Campus Governance and Web Application domain.\n"
        "In this seminar, I will walk you through our problem identification, innovative multi-role workflow, alternative strategies analyzed, system architecture, budget estimation, implementation timeline, and testing validation.\n"
        "Estimated Speaking Time: 1 minute.")

    # ==========================================
    # SLIDE 2: PROBLEM IDENTIFICATION
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Problem Identification: Challenges in Campus Facility Management")

    # 4 Problem Cards
    cards_data = [
        ("Delayed Complaint Logging", "Students notice damaged switchboards, non-functioning PCs, or broken benches but face tedious manual procedures, causing severe reporting delays.", COLOR_RED, "fa-clock"),
        ("Lack of Real-Time Tracking", "Once a complaint is reported verbally or in physical registers, students have zero visibility into who is working on it or when it will be fixed.", COLOR_AMBER, "fa-eye-slash"),
        ("Communication Gaps & No SLA", "Handoffs between admin offices, departmental heads, and technicians are completely unstructured, leading to blame-shifting and lost tickets.", COLOR_PRIMARY, "fa-comments"),
        ("Absence of Resolution Proof", "Technicians may verbally claim a repair is finished without visual proof or faculty inspection, resulting in recurring breakdowns.", COLOR_ACCENT, "fa-shield-xmark")
    ]

    for i, (title, desc, color, _) in enumerate(cards_data):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.7 + row * 2.6)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Top tag
        tag = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.2), Inches(0.25), Inches(0.25))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()

        tb = s2.shapes.add_textbox(x + Inches(0.6), y + Inches(0.15), Inches(4.8), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = f"Challenge 0{i+1}: {title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = 'Calibri'

    add_speaker_notes(s2,
        "Let us first understand the real-world problem we identified on campus.\n"
        "In colleges, students face infrastructure issues daily — open wires in corridors, projector flickers in seminar halls, network switch drops in computer labs, or plumbing leaks.\n"
        "Currently, our campus relies on verbal complaints, paper logbooks, or unofficial WhatsApp messages.\n"
        "This results in four major bottlenecks:\n"
        "1. Delayed reporting because students don't know whom to contact.\n"
        "2. Zero tracking — students have no idea if their issue was even seen.\n"
        "3. Severe communication gaps between admin and technicians.\n"
        "4. No proof of resolution, leading to false completion claims.\n"
        "Campus Connect was conceived specifically to address these institutional pain points.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 3: EXISTING SYSTEM & ITS LIMITATIONS
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Existing Systems & Root Cause Analysis")

    # Left Column: Existing Methods
    box_left = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = COLOR_CARD_BG
    box_left.line.color.rgb = COLOR_CARD_BORDER
    
    tl = box_left.text_frame
    tl.word_wrap = True
    tl.margin_left = Inches(0.3)
    tl.margin_top = Inches(0.3)
    
    p = tl.paragraphs[0]
    p.text = "CURRENT TRADITIONAL METHODS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    items_left = [
        ("Physical Register Books", "Stationed at administrative counters; prone to physical damage, illegible handwriting, and zero automated notifications."),
        ("Verbal Reporting to Peons/Staff", "High probability of human forgetfulness, lack of accountability, and no timestamped logging."),
        ("Ad-Hoc Messaging / WhatsApp", "Unstructured complaint data, unmanageable media files, no role verification, and tickets get buried in chat history."),
        ("No SLA Deadlines", "Work orders have no enforced completion target or escalation mechanism for high-priority hazards.")
    ]
    for title, desc in items_left:
        p_t = tl.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_RED
        
        p_d = tl.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_after = Pt(8)

    # Right Column: Root Cause & Impact
    box_right = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = COLOR_SECONDARY
    box_right.line.color.rgb = RGBColor(51, 65, 85)

    tr = box_right.text_frame
    tr.word_wrap = True
    tr.margin_left = Inches(0.3)
    tr.margin_top = Inches(0.3)

    p_r = tr.paragraphs[0]
    p_r.text = "INSTITUTIONAL IMPACT & CONSEQUENCES"
    p_r.font.size = Pt(14)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_ACCENT_LIGHT
    p_r.space_after = Pt(10)

    impacts = [
        ("Average Resolution Delay", "Issues take 5 to 14 days to resolve due to multi-department paper hopping."),
        ("Zero Accountable Auditing", "Campus authorities cannot evaluate technician performance or recurring appliance faults."),
        ("Safety Hazards Left Unattended", "Critical electrical sparkings or water seepage remain unaddressed for days due to lack of priority tagging."),
        ("Student Dissatisfaction", "Students lose confidence in campus administration when reported grievances vanish without updates.")
    ]
    for title, desc in impacts:
        p_t = tr.add_paragraph()
        p_t.text = f"✖ {title}"
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_AMBER
        
        p_d = tr.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_LIGHT
        p_d.space_after = Pt(8)

    add_speaker_notes(s3,
        "In this slide, we analyze why existing campus approaches fail.\n"
        "On the left, we see the three traditional mechanisms: manual registers, verbal communication, and unstructured WhatsApp groups.\n"
        "None of these offer timestamping, role segregation, or SLA enforcement.\n"
        "On the right, we show the direct consequences:\n"
        "Repairs take up to 2 weeks, safety hazards like exposed 230V wiring in corridors go unattended, and the college has zero data to review technician efficiency.\n"
        "This proves that a digitized, strictly governed platform is a necessity, not just a luxury.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 4: PROPOSED SOLUTION — CAMPUS CONNECT
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Proposed Solution: Campus Connect Architecture")

    # Banner Concept
    banner = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_ACCENT
    banner.line.fill.background()
    btf = banner.text_frame
    btf.word_wrap = True
    btf.margin_left = Inches(0.3)
    btf.margin_top = Inches(0.18)
    
    bp1 = btf.paragraphs[0]
    bp1.text = "CAMPUS CONNECT — A 5-STAGE CLOSED-LOOP RESOLUTION LIFECYCLE"
    bp1.font.size = Pt(13)
    bp1.font.bold = True
    bp1.font.color.rgb = RGBColor(255, 255, 255)

    bp2 = btf.add_paragraph()
    bp2.text = "Student (Reports Issue) → Admin (Audits & Dispatches) → Technician (Executes & Uploads Proof) → Faculty (Inspects & Verifies) → Resolution"
    bp2.font.size = Pt(12)
    bp2.font.bold = True
    bp2.font.color.rgb = COLOR_TEXT_LIGHT

    # 4 Pillar Cards
    pillars = [
        ("1. Role-Based Dashboards", "Dedicated, secure interfaces for Students, Admin, Technicians, and Faculty with strict session token validation.", COLOR_PRIMARY),
        ("2. Dual-Tier Verification", "Admin verifies inbound authenticity; Department Faculty performs final on-site inspection before ticket closure.", COLOR_TEAL),
        ("3. Visual Evidence Mandate", "Students attach fault photos/videos; Technicians MUST upload certified repair proof before submitting for QA audit.", COLOR_ACCENT),
        ("4. 7-Stage Live Progress Bar", "Students track their exact ticket milestone with transparent percentages (14% → 100% Work Done).", COLOR_EMERALD)
    ]

    for i, (title, desc, colr) in enumerate(pillars):
        x = Inches(0.8 + i * 2.98)
        y = Inches(3.05)
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.82), Inches(3.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Header band on card
        band = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(2.52), Inches(0.55))
        band.fill.solid()
        band.fill.fore_color.rgb = colr
        band.line.fill.background()
        bt = band.text_frame
        bt.word_wrap = True
        bt.margin_top = Inches(0.08)
        p = bt.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Content text
        ctb = s4.shapes.add_textbox(x + Inches(0.15), y + Inches(0.85), Inches(2.52), Inches(2.8))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        p2 = ctf.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = 'Calibri'

    add_speaker_notes(s4,
        "Here is our proposed solution: Campus Connect.\n"
        "Unlike generic ticketing tools, Campus Connect implements a strict 5-stage sequential governance model:\n"
        "Student reports -> Admin audits and dispatches -> Technician repairs -> Faculty verifies -> Work Completed.\n"
        "Notice our four core pillars:\n"
        "First, role-based separation so users only see what is relevant to them.\n"
        "Second, dual-tier verification: Admin prevents spam, and Faculty confirms actual quality.\n"
        "Third, mandatory photographic evidence — technicians cannot mark work completed without uploading a photo proof of the fix.\n"
        "Fourth, a 7-stage live progress bar giving students 100% transparency.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 5: INNOVATION & UNIQUE VALUE PROPOSITION
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Innovation: Why Campus Connect is Unique & High-Value")

    innovations = [
        ("Dual-Layer Quality Clearance", "Technicians cannot unilaterally mark tickets as resolved. The academic department faculty must physically verify the repair proof, eliminating fake resolution claims.", "Accountability"),
        ("Mandatory Photographic Audit Trail", "Both initial defect (e.g. broken bulb) and completed fix (new LED installed) require photo proofs stored permanently in the audit history.", "Transparency"),
        ("Smart Priority Hazard Escalation", "Keywords such as 'open wire', 'sparking', 'short circuit' auto-flag tickets as High Priority with urgent visual tags for immediate administrative dispatch.", "Campus Safety"),
        ("7-Stage Granular Stepper UI", "Instead of binary 'Open/Closed' states, students view 7 milestones: Submitted → Admin Verified → Tech Accepted → In Progress → Tech Completed → Faculty Verified → Completed.", "User Experience"),
        ("Technician SLA & Performance Rating", "Admin sets strict deadlines (e.g. 48-hour target) and faculty approvals automatically boost technician rating points for merit tracking.", "Data-Driven")
    ]

    for i, (title, desc, badge) in enumerate(innovations):
        y = Inches(1.65 + i * 1.05)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        # Left badge
        badge_shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(0.18), Inches(1.8), Inches(0.55))
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = COLOR_ACCENT if i % 2 == 0 else COLOR_TEAL
        badge_shape.line.fill.background()
        btf = badge_shape.text_frame
        btf.margin_top = Inches(0.08)
        bp = btf.paragraphs[0]
        bp.text = badge
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(255, 255, 255)
        bp.alignment = PP_ALIGN.CENTER

        # Title and Description
        tb = s5.shapes.add_textbox(Inches(3.0), y + Inches(0.1), Inches(9.3), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_left = tf.margin_bottom = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_speaker_notes(s5,
        "As required by the Seminar-II High Proficiency criteria, let us highlight the innovation and engineering decisions in Campus Connect.\n"
        "We didn't just build a simple CRUD form. We designed 5 key innovations:\n"
        "1. Dual-Layer Quality Clearance — ensuring faculty verification before ticket closure.\n"
        "2. Mandatory Photographic Audit Trail — eliminating disputes.\n"
        "3. Smart Priority Escalation — critical electrical hazards are instantly flagged.\n"
        "4. 7-Stage Granular Stepper UI — giving granular visibility.\n"
        "5. Technician Performance Ratings & Deadlines — introducing institutional accountability.\n"
        "Each innovation solves a specific administrative bottleneck.\n"
        "Estimated Speaking Time: 2 minutes.")

    # ==========================================
    # SLIDE 6: COMPLETE SYSTEM WORKFLOW & 7-STAGE PROGRESS
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Complete 5-Stage System Workflow & 7-Stage Progress Bar")

    # 5 Workflow Stage Boxes
    stages = [
        ("Step 1: Student", "Files complaint with category, location, priority & photo evidence.", "Complaint Submitted\n(Stage 1 • 14%)", COLOR_PRIMARY),
        ("Step 2: Admin", "Audits inbound report, assigns specialized technician & deadline date.", "Approved by Admin\n(Stage 2 • 28%)", COLOR_ACCENT),
        ("Step 3: Tech Action", "Technician accepts work order and initiates physical repairs.", "Work in Progress\n(Stage 4 • 57%)", COLOR_AMBER),
        ("Step 4: Tech Proof", "Technician finishes repair, uploads completion proof photo & remarks.", "Tech Completed\n(Stage 5 • 71%)", COLOR_TEAL),
        ("Step 5: Faculty QA", "Faculty inspects fix, audits proof, and confirms closure.", "Completed ✅\n(Stage 7 • 100%)", COLOR_EMERALD)
    ]

    for i, (title, desc, status_txt, colr) in enumerate(stages):
        x = Inches(0.8 + i * 2.38)
        y = Inches(1.65)
        
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.25), Inches(3.4))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = colr
        box.line.width = Pt(2)

        # Header tag
        tag = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.1), Inches(2.05), Inches(0.5))
        tag.fill.solid()
        tag.fill.fore_color.rgb = colr
        tag.line.fill.background()
        tt = tag.text_frame
        tt.margin_top = Inches(0.06)
        p = tt.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Description
        tb = s6.shapes.add_textbox(x + Inches(0.1), y + Inches(0.7), Inches(2.05), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXT_DARK

        # Status badge inside box
        s_badge = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(2.4), Inches(2.05), Inches(0.8))
        s_badge.fill.solid()
        s_badge.fill.fore_color.rgb = COLOR_BG_LIGHT
        s_badge.line.color.rgb = colr
        st = s_badge.text_frame
        st.margin_top = Inches(0.06)
        p = st.paragraphs[0]
        p.text = status_txt
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = colr
        p.alignment = PP_ALIGN.CENTER

    # Bottom Banner: 7-Stage Progress Bar Representation
    p_banner = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.75))
    p_banner.fill.solid()
    p_banner.fill.fore_color.rgb = COLOR_SECONDARY
    p_banner.line.color.rgb = RGBColor(51, 65, 85)

    pbtf = p_banner.text_frame
    pbtf.word_wrap = True
    pbtf.margin_left = Inches(0.3)
    pbtf.margin_top = Inches(0.18)

    p = pbtf.paragraphs[0]
    p.text = "7-STAGE STUDENT PROGRESS CHECKLIST (LIVE STEPPER)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_LIGHT

    p2 = pbtf.add_paragraph()
    p2.text = "✓ 1. Complaint Submitted (14%)  →  ✓ 2. Admin Verified (28%)  →  ✓ 3. Tech Accepted (43%)  →  ✓ 4. Work in Progress (57%)  →  ✓ 5. Tech Completed (71%)  →  ✓ 6. Faculty Verified (86%)  →  ✓ 7. Completed ✅ (100% • Work Done)"
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_EMERALD
    p2.space_after = Pt(4)

    p3 = pbtf.add_paragraph()
    p3.text = "• Handles Rejection Pathways: Admin Rejection or Technician Decline immediately updates status and halts progression with recorded rationale."
    p3.font.size = Pt(10.5)
    p3.font.color.rgb = COLOR_TEXT_LIGHT

    add_speaker_notes(s6,
        "This slide illustrates the complete, unbroken workflow of Campus Connect.\n"
        "Step 1: Student files a ticket with location and evidence — status is 'Complaint Submitted' at 14%.\n"
        "Step 2: Admin reviews and assigns a specialized technician with a deadline — status becomes 'Approved by Admin' at 28%.\n"
        "Step 3: The assigned technician receives the work order and accepts — moving to 'Work in Progress' at 57%.\n"
        "Step 4: Technician completes the physical fix and uploads a proof photo — moving to 'Tech Completed' at 71%.\n"
        "Step 5: The department faculty inspects the fix and gives final QA approval — reaching 'Completed' at 100%.\n"
        "If rejected at any stage, the system records the reason and notifies the student.\n"
        "Estimated Speaking Time: 2 minutes.")

    # ==========================================
    # SLIDE 7: SYSTEM ARCHITECTURE & TECHNOLOGY STACK
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "System Architecture & 3-Tier Technology Stack")

    # 3 Architecture Tier Cards
    tiers = [
        ("Presentation Layer (Frontend)", [
            "HTML5 & Semantic Structure",
            "CSS3 & Custom Glassmorphism Theme",
            "Bootstrap 5 (Responsive Grid & Modals)",
            "Vanilla JavaScript (ES6+ Asynchronous Logic)",
            "FontAwesome 6 (Icons & Visual Cues)"
        ], COLOR_ACCENT),
        ("Application & Logic Layer (Backend)", [
            "PHP 8.x (Server-side Session & API Handler)",
            "Role-Based Authentication Engine",
            "State Transition & Workflow Controller",
            "File Upload & Base64 Media Validator",
            "RESTful JSON Data Exchange"
        ], COLOR_TEAL),
        ("Data Persistence Layer (Database)", [
            "MySQL 8.0 Relational Database",
            "Normalized Entities (3NF)",
            "Foreign Key Integrity Constraints",
            "Timestamped Resolution Audit Logs",
            "Optimized Indexing for Fast Querying"
        ], COLOR_PRIMARY)
    ]

    for i, (title, items, colr) in enumerate(tiers):
        x = Inches(0.8 + i * 3.95)
        y = Inches(1.7)
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Header banner
        band = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.65))
        band.fill.solid()
        band.fill.fore_color.rgb = colr
        band.line.fill.background()
        bt = band.text_frame
        bt.margin_top = Inches(0.1)
        p = bt.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # List
        tb = s7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.95), Inches(3.4), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.text = f"✔ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.font.name = 'Calibri'
            p.space_after = Pt(10)

    add_speaker_notes(s7,
        "Here we present the 3-Tier System Architecture and Technology Stack.\n"
        "Tier 1 is our Presentation Layer: We used HTML5, CSS3, Bootstrap 5, and modern JavaScript to create responsive, accessible dashboards for desktop and mobile.\n"
        "Tier 2 is our Application Layer: Built with PHP, handling session management, role verification, state transitions, and file handling.\n"
        "Tier 3 is our Persistence Layer: Powered by MySQL, with fully normalized tables ensuring data integrity and fast retrieval.\n"
        "This 3-tier structure makes the application modular, easy to maintain, and highly scalable across departments.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 8: MAJOR MODULES & DATABASE DESIGN
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Major Functional Modules & Database Design")

    # Left: 4 Modules
    m_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.2))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = COLOR_CARD_BG
    m_box.line.color.rgb = COLOR_CARD_BORDER
    
    mtf = m_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.3)
    mtf.margin_top = Inches(0.25)
    
    p = mtf.paragraphs[0]
    p.text = "CORE FUNCTIONAL MODULES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    modules = [
        ("Student Module", "Issue registration, priority selection, media evidence upload, 7-stage live progress tracking, and resolution history."),
        ("Admin Console", "Inbound ticket audit, specialized technician dispatching, deadline assignment, technician account management, and analytical reports."),
        ("Technician Module", "Work order acceptance/rejection, task deadline monitor, resolution photo proof upload, and remarks submission."),
        ("Faculty QA Panel", "Department-level work order audit, photographic proof inspection, and final quality completion clearance.")
    ]
    for title, desc in modules:
        p1 = mtf.add_paragraph()
        p1.text = f"• {title}"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT
        
        p2 = mtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_after = Pt(6)

    # Right: Database Entities Table
    d_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
    d_box.fill.solid()
    d_box.fill.fore_color.rgb = COLOR_CARD_BG
    d_box.line.color.rgb = COLOR_CARD_BORDER

    dtf = d_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = Inches(0.3)
    dtf.margin_top = Inches(0.25)

    p = dtf.paragraphs[0]
    p.text = "DATABASE SCHEMA (RELATIONAL ENTITIES)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    entities = [
        ("tbl_users / students", "gr_no (PK), name, password_hash, dept, role, status, avatar"),
        ("tbl_complaints", "complaint_id (PK), title, category, description, location, priority, reported_by, reported_by_gr, status, stage, admin_status, technician_status, work_status, faculty_status, tech_id, deadline, image_url, video_url, proof_img, remark, qa_verified, qa_feedback"),
        ("tbl_technicians", "tech_id (PK), name, dept, experience, rating, active_status"),
        ("tbl_audit_logs", "log_id (PK), complaint_id (FK), status_tag, note, action_by, timestamp"),
        ("tbl_departments", "dept_id (PK), dept_name, faculty_advisor_id")
    ]
    for ent, fields in entities:
        p1 = dtf.add_paragraph()
        p1.text = f"🗄 {ent}"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEAL
        
        p2 = dtf.add_paragraph()
        p2.text = fields
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_after = Pt(4)

    add_speaker_notes(s8,
        "In this slide, we present the functional modules and database design.\n"
        "The system consists of four distinct modules: Student, Admin, Technician, and Faculty.\n"
        "On the database side, we have structured the schema with 5 key relational entities:\n"
        "1. tbl_users storing student credentials and roles.\n"
        "2. tbl_complaints storing the primary ticket state, stage index, assigned technician, and proof image URLs.\n"
        "3. tbl_technicians maintaining department specializations and rating metrics.\n"
        "4. tbl_audit_logs recording every status transition timestamp for a complete historical trail.\n"
        "5. tbl_departments mapping academic faculties.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 9: ALTERNATIVE SOLUTIONS & COMPARATIVE ANALYSIS
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Alternative Solutions Considered & Comparative Analysis")

    # Table of comparison
    # Rows: Metric, Manual Register, Google Forms, WhatsApp Groups, Campus Connect
    table_shape = s9.shapes.add_table(6, 5, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.2))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.3)
    table.columns[4].width = Inches(2.6)

    headers = ["Evaluation Metric", "Manual Register", "Google Forms + Sheets", "WhatsApp Groups", "Campus Connect (Proposed)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY if col_idx < 4 else COLOR_ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    data = [
        ("Automated Multi-Role Routing", "❌ None (Manual Paper)", "❌ None (Flat Spreadsheet)", "❌ None (Unstructured)", "✔ Full 5-Stage Automation"),
        ("Live Progress Tracking", "❌ Zero Visibility", "⚠ Partial (Sheet Status)", "❌ Lost in Chat Feed", "✔ 7-Stage Live Stepper"),
        ("Photo Resolution Proof", "❌ No Media Support", "⚠ Link Only", "⚠ Cluttered Chat Media", "✔ Mandatory Upload & QA"),
        ("Faculty Quality Verification", "❌ Unaudited", "❌ No Verification Step", "❌ No Verification Step", "✔ Enforced Dual-Tier QA"),
        ("Institutional Audit & SLA", "❌ Vulnerable to Loss", "⚠ Basic Filter", "❌ Zero Audit Trail", "✔ Full History & CSV Export")
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            if col_idx == 4:
                cell.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light blue accent
            else:
                cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else COLOR_BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if col_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_ACCENT
                elif col_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PRIMARY
                else:
                    p.font.color.rgb = COLOR_TEXT_DARK
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    # Bottom conclusion bar
    c_bar = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.9))
    c_bar.fill.solid()
    c_bar.fill.fore_color.rgb = COLOR_SECONDARY
    c_bar.line.fill.background()
    ctf = c_bar.text_frame
    ctf.margin_left = Inches(0.3)
    ctf.margin_top = Inches(0.12)
    p = ctf.paragraphs[0]
    p.text = "CONCLUSION ON STRATEGY SELECTION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p2 = ctf.add_paragraph()
    p2.text = "Campus Connect is selected because it is the only solution that integrates role segregation, SLA enforcement, mandatory photo proof, and faculty QA verification in a single lightweight, zero-license-cost architecture."
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_LIGHT

    add_speaker_notes(s9,
        "A critical requirement of the Seminar-II rubric is analyzing alternative solutions.\n"
        "We systematically evaluated 3 alternative approaches against Campus Connect:\n"
        "1. Manual Paper Registers: Zero automation, zero tracking, high risk of physical record loss.\n"
        "2. Google Forms + Sheets: While easy to create, it lacks role-based access, cannot enforce technician photo uploads, and has no closed-loop verification.\n"
        "3. WhatsApp Groups: Completely unstructured, media gets mixed with personal chats, and there is no accountability.\n"
        "As summarized in our comparison matrix, Campus Connect is the only platform providing end-to-end multi-role routing, photographic proof, and faculty verification.\n"
        "Estimated Speaking Time: 2 minutes.")

    # ==========================================
    # SLIDE 10: PROJECT BUDGET & COST ANALYSIS
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Project Budget & Cost Analysis (Development vs. Deployment)")

    # Left: Development Cost (Open Source & Zero License)
    dev_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.2))
    dev_box.fill.solid()
    dev_box.fill.fore_color.rgb = COLOR_CARD_BG
    dev_box.line.color.rgb = COLOR_CARD_BORDER
    dev_box.line.width = Pt(1.5)

    dtf = dev_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = Inches(0.3)
    dtf.margin_top = Inches(0.25)

    p = dtf.paragraphs[0]
    p.text = "A. DEVELOPMENT PHASE COST (STUDENT PROJECT)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    dev_items = [
        ("Frontend Technologies", "HTML5, CSS3, JavaScript, Bootstrap 5", "₹0 (Open Source)"),
        ("Backend & Database", "PHP 8.x, MySQL Community Edition", "₹0 (Open Source)"),
        ("IDE & Development Tools", "VS Code, Git, Chrome DevTools", "₹0 (Free / Open Source)"),
        ("Hardware (Existing)", "College / Student PC / Laptop", "₹0 (Existing Infrastructure)"),
        ("Local Testing Server", "XAMPP / Apache Local Environment", "₹0 (Free Software)"),
        ("Total Development Cost", "Leveraged 100% Free & Open Source Stack", "₹0 (ZERO COST)")
    ]

    for title, desc, cost in dev_items:
        p = dtf.add_paragraph()
        p.text = f"• {title}: {desc} — "
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = 'Calibri'
        
        # Add bold cost
        run = p.add_run()
        run.text = cost
        run.font.bold = True
        run.font.color.rgb = COLOR_EMERALD if "₹0" in cost else COLOR_ACCENT

    # Right: Institutional Deployment Cost Estimate
    dep_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.2))
    dep_box.fill.solid()
    dep_box.fill.fore_color.rgb = COLOR_CARD_BG
    dep_box.line.color.rgb = COLOR_CARD_BORDER
    dep_box.line.width = Pt(1.5)

    deptf = dep_box.text_frame
    deptf.word_wrap = True
    deptf.margin_left = Inches(0.3)
    deptf.margin_top = Inches(0.25)

    p = deptf.paragraphs[0]
    p.text = "B. ESTIMATED INSTITUTIONAL DEPLOYMENT (ANNUAL)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    dep_items = [
        ("Campus Domain Registration", ".edu.in or .ac.in sub-domain", "₹800 – ₹1,200 / year"),
        ("Cloud / Shared Web Hosting", "Linux CPanel / VPS with PHP & MySQL", "₹2,500 – ₹4,000 / year"),
        ("SSL Security Certificate", "Let's Encrypt TLS Certificate", "₹0 (Free Open Source)"),
        ("Campus Local Server Option", "Host on College Intranet Server", "₹0 (Using College LAN)"),
        ("Annual Maintenance (AMC)", "In-house Student / IT Staff Maintenance", "₹0 (Internal Support)"),
        ("Total Estimated Deployment", "Full institutional rollout range", "₹3,300 – ₹5,200 / year")
    ]

    for title, desc, cost in dep_items:
        p = deptf.add_paragraph()
        p.text = f"• {title}: {desc} — "
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = 'Calibri'
        
        run = p.add_run()
        run.text = cost
        run.font.bold = True
        run.font.color.rgb = COLOR_ACCENT

    # Budget Summary Note
    b_note = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.95))
    b_note.fill.solid()
    b_note.fill.fore_color.rgb = COLOR_SECONDARY
    b_note.line.fill.background()
    bntf = b_note.text_frame
    bntf.margin_left = Inches(0.3)
    bntf.margin_top = Inches(0.12)
    p = bntf.paragraphs[0]
    p.text = "FINANCIAL FEASIBILITY HIGHLIGHT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    p2 = bntf.add_paragraph()
    p2.text = "By deliberately choosing an open-source FOSS technology stack (PHP, MySQL, Bootstrap), our development cost is ₹0. For college deployment, the system can run on the existing campus intranet server for ₹0 additional expenditure."
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_LIGHT

    add_speaker_notes(s10,
        "Cost and budget planning are essential parts of our Seminar-II evaluation.\n"
        "We have categorized our budget into two distinct sections:\n"
        "1. Development Cost: Exactly ₹0. We intentionally selected open-source technologies — PHP, MySQL, Bootstrap, and VS Code — eliminating any proprietary licensing fees.\n"
        "2. Deployment Cost: If deployed on commercial cloud hosting with a dedicated domain, the estimated annual cost is between ₹3,300 to ₹5,200.\n"
        "Alternatively, the college can host it on the existing campus local intranet server at zero additional cost.\n"
        "This makes Campus Connect economically feasible for any educational institution.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 11: DEVELOPMENT PLAN & TIMELINE (GANTT)
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "Development Plan & Implementation Timeline")

    # 5 Major Phase Blocks
    phases = [
        ("Phase 1: Research & SRS", "Weeks 1 – 3", [
            "Problem identification across campus",
            "Requirement gathering from students & staff",
            "Feasibility study & alternative analysis",
            "Finalizing Software Requirement Specification"
        ], COLOR_PRIMARY),
        ("Phase 2: UI/UX & DB Design", "Weeks 4 – 6", [
            "Entity-Relationship (ER) modeling in 3NF",
            "Designing database tables in MySQL",
            "Figma wireframing & responsive layouts",
            "Design system setup with Bootstrap 5"
        ], COLOR_ACCENT),
        ("Phase 3: Core Implementation", "Weeks 7 – 10", [
            "Role authentication & session management",
            "Student complaint submission module",
            "Admin dispatch & technician routing",
            "Technician proof upload & Faculty QA logic"
        ], COLOR_TEAL),
        ("Phase 4: Testing & Bug Fixing", "Weeks 11 – 13", [
            "Unit testing of PHP endpoints & forms",
            "End-to-end workflow validation tests",
            "Security testing for unauthorized role access",
            "Cross-browser & mobile responsiveness"
        ], COLOR_AMBER),
        ("Phase 5: Deployment & Docs", "Weeks 14 – 16", [
            "Intranet server deployment & DB migration",
            "User acceptance testing (UAT)",
            "Seminar-II presentation & report preparation",
            "Project manual & documentation finalized"
        ], COLOR_EMERALD)
    ]

    for i, (p_title, p_time, tasks, colr) in enumerate(phases):
        x = Inches(0.8 + i * 2.38)
        y = Inches(1.7)
        card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.25), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Header tag
        tag = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.1), Inches(2.05), Inches(0.65))
        tag.fill.solid()
        tag.fill.fore_color.rgb = colr
        tag.line.fill.background()
        tt = tag.text_frame
        tt.margin_top = Inches(0.08)
        p = tt.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tt.add_paragraph()
        p_sub.text = p_time
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = COLOR_BG_LIGHT
        p_sub.alignment = PP_ALIGN.CENTER

        # Tasks
        tb = s11.shapes.add_textbox(x + Inches(0.1), y + Inches(0.85), Inches(2.05), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for t in tasks:
            p = tf.add_paragraph()
            p.text = f"• {t}"
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.font.name = 'Calibri'
            p.space_after = Pt(8)

    add_speaker_notes(s11,
        "Here is our 16-week structured implementation timeline.\n"
        "We followed an agile engineering lifecycle across 5 distinct phases:\n"
        "Phase 1 (Weeks 1-3): Requirement gathering and SRS finalization.\n"
        "Phase 2 (Weeks 4-6): Database schema design in 3NF and UI wireframing.\n"
        "Phase 3 (Weeks 7-10): Core development of all 4 role portals and workflow state transitions.\n"
        "Phase 4 (Weeks 11-13): Comprehensive testing including security and role-boundary validation.\n"
        "Phase 5 (Weeks 14-16): Deployment, User Acceptance Testing, and technical documentation.\n"
        "This systematic plan ensures timely completion and high quality.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 12: FEASIBILITY ANALYSIS
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_header(s12, "Feasibility Analysis: Technical, Economic, Operational & Scalability")

    feasibility_data = [
        ("Technical Feasibility", "HIGH", [
            "PHP and MySQL are industry-standard, mature technologies supported natively across web servers.",
            "Modern browsers fully support HTML5 media capture, JavaScript async requests, and responsive CSS grids.",
            "No specialized, complex third-party hardware or proprietary software drivers are required."
        ], COLOR_ACCENT),
        ("Economic Feasibility", "HIGH", [
            "Utilizes a 100% Free and Open Source Software (FOSS) stack, requiring ₹0 licensing investment.",
            "Can be hosted entirely on existing college server infrastructure or low-cost cloud hosting (₹3,300/yr).",
            "Significantly reduces paper, manual register stationery, and administrative labor costs."
        ], COLOR_EMERALD),
        ("Operational Feasibility", "HIGH", [
            "Intuitive UI requires zero training for students and faculty members.",
            "Technicians receive simplified action buttons: 'Accept', 'Decline', and 'Upload Proof'.",
            "Clear visual feedback and transparent progress bars encourage campus-wide adoption."
        ], COLOR_TEAL),
        ("Scalability Feasibility", "HIGH", [
            "Modular database architecture easily scales from single building to multi-campus institutions.",
            "Can easily accommodate new departments, additional technicians, and specialized facility categories.",
            "Designed to support future AI/ML automated classification and SMS/Email notification APIs."
        ], COLOR_PRIMARY)
    ]

    for i, (title, score, points, colr) in enumerate(feasibility_data):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.7 + row * 2.6)

        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        # Header tag
        tag = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.18), Inches(3.2), Inches(0.45))
        tag.fill.solid()
        tag.fill.fore_color.rgb = colr
        tag.line.fill.background()
        tt = tag.text_frame
        tt.margin_top = Inches(0.06)
        p = tt.paragraphs[0]
        p.text = f"{title} [{score}]"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Points
        tb = s12.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(5.2), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        for pt in points:
            p = tf.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(10.5)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.font.name = 'Calibri'
            p.space_after = Pt(4)

    add_speaker_notes(s12,
        "We conducted a 4-dimensional feasibility study:\n"
        "1. Technical Feasibility: PHP and MySQL are proven, robust, and supported everywhere.\n"
        "2. Economic Feasibility: Zero software licensing costs; runs on existing college hardware.\n"
        "3. Operational Feasibility: Extremely simple role-based user interfaces with minimal learning curve.\n"
        "4. Scalability: The database schema can effortlessly expand to support multiple campuses and thousands of concurrent users.\n"
        "All four dimensions confirm that Campus Connect is 100% viable for real-world deployment.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 13: TESTING, VALIDATION & TEST CASES
    # ==========================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13)
    add_header(s13, "Testing, Validation & Quality Assurance")

    # Table of Test Cases
    t_shape = s13.shapes.add_table(7, 5, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.2))
    table = t_shape.table

    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.2)
    table.columns[4].width = Inches(1.5)

    headers = ["Test ID", "Test Scenario", "Input / Action", "Expected Outcome", "Status"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    test_cases = [
        ("TC-01", "Student Complaint Submission", "Valid title, category, location, priority & image", "Ticket saved, status set to 'Complaint Submitted' (Stage 1 • 14%)", "PASSED ✔"),
        ("TC-02", "Admin Inbound Audit & Dispatch", "Select technician & assign deadline date", "Ticket status updated to 'Approved by Admin' (Stage 2 • 28%)", "PASSED ✔"),
        ("TC-03", "Technician Task Acceptance", "Click 'Accept Complaint' button", "Status updated to 'Work in Progress' (Stage 4 • 57%)", "PASSED ✔"),
        ("TC-04", "Technician Work Completion", "Upload base64 proof photo and resolution comments", "Status becomes 'Work Completed by Technician' (Stage 5 • 71%)", "PASSED ✔"),
        ("TC-05", "Faculty QA Final Verification", "Inspect photo proof and click 'Verify & Approve'", "Status becomes 'Completed ✅' (Stage 7 • 100% Work Done)", "PASSED ✔"),
        ("TC-06", "Role Security & Unauthorized Access", "Student attempting to access Admin / Faculty panel", "Access blocked, session redirected to authorized dashboard", "PASSED ✔")
    ]

    for row_idx, row_data in enumerate(test_cases, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else COLOR_BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                if col_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_EMERALD
                    p.alignment = PP_ALIGN.CENTER
                elif col_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PRIMARY
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = COLOR_TEXT_DARK

    # Summary bar
    s_bar = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.95))
    s_bar.fill.solid()
    s_bar.fill.fore_color.rgb = COLOR_SECONDARY
    s_bar.line.fill.background()
    sbtf = s_bar.text_frame
    sbtf.margin_left = Inches(0.3)
    sbtf.margin_top = Inches(0.12)
    p = sbtf.paragraphs[0]
    p.text = "TESTING & AUTOMATED SUITE RESULTS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p2 = sbtf.add_paragraph()
    p2.text = "Our automated test suite validated 9 distinct end-to-end user journeys including student submission, admin dispatch, technician proof validation, faculty QA approvals, and rejection paths with 100% passing test results."
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_LIGHT

    add_speaker_notes(s13,
        "To ensure robust software quality, we created comprehensive test cases.\n"
        "We tested all normal flows as well as edge cases:\n"
        "Test cases 1 through 5 validated the sequential lifecycle transitions from Stage 1 (14%) to Stage 7 (100%).\n"
        "Test case 6 tested role-based security boundaries, confirming that students cannot access admin or faculty endpoints.\n"
        "We also ran automated validation scripts that verified all 9 test journeys with zero errors.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 14: EXPECTED OUTCOMES & FUTURE SCOPE
    # ==========================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14)
    add_header(s14, "Expected Outcomes & Future Enhancements")

    # Left: Expected Outcomes
    o_box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.2))
    o_box.fill.solid()
    o_box.fill.fore_color.rgb = COLOR_CARD_BG
    o_box.line.color.rgb = COLOR_CARD_BORDER
    o_box.line.width = Pt(1.5)

    otf = o_box.text_frame
    otf.word_wrap = True
    otf.margin_left = Inches(0.3)
    otf.margin_top = Inches(0.25)

    p = otf.paragraphs[0]
    p.text = "MEASURABLE PROJECT OUTCOMES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    outcomes = [
        ("90% Faster Issue Turnaround", "Reduces average campus repair cycle from 7–14 days down to 24–48 hours through automated dispatch."),
        ("100% Real-Time Transparency", "Eliminates student uncertainty with continuous 7-stage status updates and timestamps."),
        ("Zero Paper Waste & Complete Audit Trail", "Replaces physical logbooks with secure, centralized, and downloadable digital complaint histories."),
        ("Enforced Quality Assurance", "Mandatory technician repair photos and faculty verification prevent premature or false ticket closures.")
    ]
    for title, desc in outcomes:
        p1 = otf.add_paragraph()
        p1.text = f"✔ {title}"
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_EMERALD
        
        p2 = otf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_after = Pt(6)

    # Right: Future Scope
    f_box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = COLOR_CARD_BG
    f_box.line.color.rgb = COLOR_CARD_BORDER
    f_box.line.width = Pt(1.5)

    ftf = f_box.text_frame
    ftf.word_wrap = True
    ftf.margin_left = Inches(0.3)
    ftf.margin_top = Inches(0.25)

    p = ftf.paragraphs[0]
    p.text = "PLANNED FUTURE ENHANCEMENTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    futures = [
        ("AI/ML Natural Language Categorization", "Integrate Naive Bayes / NLP models to automatically predict complaint category and priority from student descriptions."),
        ("QR-Code Smart Facility Tagging", "Affix QR codes to classroom projectors, switchboards, and lab PCs for instant 1-scan complaint filing."),
        ("Automated SMS & WhatsApp Alerts", "Integrate Twilio / WhatsApp Business API for instant push alerts to technicians upon ticket dispatch."),
        ("Computer Vision Damage Detection", "Use image classification models (CNN/TensorFlow) to auto-verify appliance damage in uploaded photos.")
    ]
    for title, desc in futures:
        p1 = ftf.add_paragraph()
        p1.text = f"🚀 {title}"
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT
        
        p2 = ftf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_after = Pt(6)

    add_speaker_notes(s14,
        "Looking at the outcomes and roadmap:\n"
        "The direct outcomes of Campus Connect are a 90% reduction in resolution delays, 100% tracking transparency, elimination of paper logbooks, and verified repair quality.\n"
        "For future enhancements, we have identified 4 realistic technical additions:\n"
        "1. Machine Learning NLP for automated category and priority prediction.\n"
        "2. QR-code tagging on campus appliances for instant reporting.\n"
        "3. Automated SMS/WhatsApp notifications for technicians.\n"
        "4. Computer vision models to detect physical damage in uploaded photos.\n"
        "Estimated Speaking Time: 1.5 minutes.")

    # ==========================================
    # SLIDE 15: CONCLUSION & RUBRIC ALIGNMENT
    # ==========================================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_background(s15, COLOR_PRIMARY)

    # Decorative accent bar
    accent_bar = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()

    # Content box
    c_box = s15.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11.0), Inches(5.5))
    ctf = c_box.text_frame
    ctf.word_wrap = True

    p0 = ctf.paragraphs[0]
    p0.text = "SEMINAR-II CONCLUSION & RUBRIC SUMMARY"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_LIGHT
    p0.space_after = Pt(12)

    p1 = ctf.add_paragraph()
    p1.text = "Campus Connect Delivers a Complete, Feasible & Innovative Solution"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.space_after = Pt(20)

    rubric_points = [
        ("INNOVATION", "Dual-tier verification, mandatory photo proof, smart priority hazard tags, and 7-stage live progress tracking."),
        ("REALISTIC PLAN", "Structured 16-week agile timeline covering SRS, 3NF database design, core module coding, and validation."),
        ("BUDGET ANALYSIS", "100% open-source stack (₹0 dev cost) with clear institutional deployment estimate (₹3,300–₹5,200/yr)."),
        ("ALTERNATIVE STRATEGIES", "Systematically compared against Manual Registers, Google Forms, and WhatsApp Groups with clear justification.")
    ]

    for title, desc in rubric_points:
        p = ctf.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_EMERALD
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_LIGHT
        p.space_after = Pt(8)

    p_end = ctf.add_paragraph()
    p_end.text = "\nThank You! Open for Questions & Technical Feedback."
    p_end.font.size = Pt(18)
    p_end.font.bold = True
    p_end.font.color.rgb = COLOR_ACCENT_LIGHT
    p_end.alignment = PP_ALIGN.CENTER

    add_speaker_notes(s15,
        "In conclusion, Campus Connect fulfills all four core requirements of the Seminar-II High Proficiency rubric:\n"
        "- Innovation: Through our dual-tier verification, mandatory photographic proof, and 7-stage live tracking.\n"
        "- Realistic Plan: Outlined through our 16-week phased development timeline.\n"
        "- Budget: Documenting a ₹0 open-source development cost and affordable deployment model.\n"
        "- Alternatives: Demonstrated through our comparative analysis against manual registers and Google Forms.\n"
        "The project is fully feasible, functional, and ready for campus deployment.\n"
        "Thank you respected evaluators. I am now open to your questions and feedback.\n"
        "Estimated Speaking Time: 1 minute.")

    output_path = 'c:/Users/Darshan/Campus---Connect/Campus_Connect_Seminar_II_Presentation.pptx'
    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_presentation()
